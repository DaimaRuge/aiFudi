# File Extractor
#文件提取器

import hashlib
from pathlib import Path
from typing import Dict, Optional, Any
from ..core.config import settings
import fitz  # PyMuPDF
import docx
import pptx
from datetime import datetime


class FileExtractor:
    """
    文件提取器
    
    支持多种文档格式的元数据提取和内容解析
    """
    
    # 支持的扩展名
    SUPPORTED_EXTENSIONS = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".html": "text/html",
    }
    
    def __init__(self):
        pass
    
    def is_supported(self, extension: str) -> bool:
        """检查是否支持该扩展名"""
        return extension.lower() in self.SUPPORTED_EXTENSIONS
    
    def extract(self, file_path: str) -> Dict[str, Any]:
        """
        提取文件元数据
        
        Args:
            file_path: 文件路径
            
        Returns:
            Dict: 元数据字典
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        
        extractor = self._get_extractor(extension)
        
        if extractor:
            return extractor(path)
        else:
            return self._extract_basic(path)
    
    def _get_extractor(self, extension: str):
        """获取对应的提取器"""
        extractors = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".txt": self._extract_text,
            ".md": self._extract_markdown,
        }
        return extractors.get(extension)
    
    def _extract_basic(self, path: Path) -> Dict[str, Any]:
        """基础提取 (所有文件通用)"""
        return {
            "title": path.stem,
            "file_size": path.stat().st_size if path.exists() else 0,
        }
    
    def _extract_pdf(self, path: Path) -> Dict[str, Any]:
        """提取 PDF 元数据"""
        try:
            doc = fitz.open(str(path))
            
            metadata = doc.metadata
            
            # 提取文本内容
            text_content = ""
            for page in doc:
                text_content += page.get_text()
            
            result = {
                "title": metadata.get("title") or path.stem,
                "authors": [metadata.get("author", "")] if metadata.get("author") else [],
                "abstract": metadata.get("subject", ""),
                "keywords": metadata.get("keywords", "").split(",") if metadata.get("keywords") else [],
                "content_text": text_content.strip()[:50000],  # 
            }
            
            doc.close()
限制长度            return result
            
        except Exception as e:
            print(f"PDF 提取失败: {e}")
            return self._extract_basic(path)
    
    def _extract_docx(self, path: Path) -> Dict[str, Any]:
        """提取 DOCX 元数据"""
        try:
            doc = docx.Document(str(path))
            
            # 提取段落文本
            paragraphs = [p.text for p in doc.paragraphs]
            text_content = "\n".join(paragraphs)
            
            # 提取属性
            core_props = doc.core_properties
            
            result = {
                "title": core_props.title or path.stem,
                "authors": [core_props.author] if core_props.author else [],
                "abstract": core_props.comments or "",
                "keywords": core_props.keywords.split(",") if core_props.keywords else [],
                "content_text": text_content.strip()[:50000],
            }
            
            return result
            
        except Exception as e:
            print(f"DOCX 提取失败: {e}")
            return self._extract_basic(path)
    
    def _extract_pptx(self, path: Path) -> Dict[str, Any]:
        """提取 PPTX 元数据"""
        try:
            prs = pptx.Presentation(str(path))
            
            # 提取幻灯片文本
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            
            text_content = "\n".join(slides_text)
            
            core_props = prs.core_properties
            
            result = {
                "title": core_props.title or path.stem,
                "authors": [core_props.author] if core_props.author else [],
                "abstract": core_props.comments or "",
                "content_text": text_content.strip()[:50000],
            }
            
            return result
            
        except Exception as e:
            print(f"PPTX 提取失败: {e}")
            return self._extract_basic(path)
    
    def _extract_text(self, path: Path) -> Dict[str, Any]:
        """提取纯文本"""
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            
            # 尝试从第一行提取标题
            lines = content.strip().split("\n")
            title = lines[0][:100] if lines else path.stem
            
            return {
                "title": title,
                "content_text": content[:50000],
            }
            
        except Exception as e:
            print(f"TXT 提取失败: {e}")
            return self._extract_basic(path)
    
    def _extract_markdown(self, path: Path) -> Dict[str, Any]:
        """提取 Markdown"""
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            
            # 提取标题 (第一个 # 开头的内容)
            lines = content.strip().split("\n")
            title = path.stem
            for line in lines:
                if line.startswith("# "):
                    title = line[2:].strip()
                    break
            
            return {
                "title": title,
                "content_text": content[:50000],
            }
            
        except Exception as e:
            print(f"Markdown 提取失败: {e}")
            return self._extract_basic(path)
    
    @staticmethod
    def compute_hash(file_path: str) -> str:
        """计算文件 SHA-256 哈希"""
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()
    
    @staticmethod
    def get_file_info(file_path: str) -> Dict[str, Any]:
        """获取文件基本信息"""
        path = Path(file_path)
        return {
            "file_name": path.name,
            "file_size": path.stat().st_size if path.exists() else 0,
            "extension": path.suffix.lower(),
            "created_at": datetime.fromtimestamp(path.stat().st_ctime),
            "modified_at": datetime.fromtimestamp(path.stat().st_mtime),
        }
